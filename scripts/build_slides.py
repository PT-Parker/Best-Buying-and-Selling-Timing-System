import os
import argparse
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

import yaml


def _load_config(path: str) -> dict:
    with open(path, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def _resolve_symbols(args) -> List[str]:
    if args.symbols:
        return args.symbols
    if args.config and os.path.exists(args.config):
        cfg = _load_config(args.config)
        syms = cfg.get('symbols') or []
        return list(syms)
    return []


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_symbol_slide(prs: Presentation, symbol: str, img_path: str, notes: List[str]):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"{symbol} 回測摘要"

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width=width)
    else:
        tx = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf = tx.text_frame
        tf.text = f"找不到圖檔: {img_path}"
        tf.paragraphs[0].font.size = Pt(16)

    # Notes box
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
    tf2 = tx2.text_frame
    tf2.text = "\n".join(notes)
    tf2.paragraphs[0].font.size = Pt(14)


def build(args) -> int:
    symbols = _resolve_symbols(args)
    if not symbols:
        print("No symbols provided (use --symbols or --config)")
        return 2

    os.makedirs(os.path.dirname(args.out) or '.', exist_ok=True)

    prs = Presentation()
    add_title_slide(prs, args.title, "由回測輸出自動產生")

    for sym in symbols:
        img = os.path.join('backtest_out', f'{sym}_equity.png')
        notes = [
            f"圖表：{os.path.relpath(img)}",
            f"交易紀錄：backtest_out/{sym}_trades.csv",
        ]
        add_symbol_slide(prs, sym, img, notes)

    prs.save(args.out)
    print(f"Slides saved to {args.out}")
    return 0


def parse_args():
    p = argparse.ArgumentParser(description='Build simple PPTX slides from backtest outputs')
    p.add_argument('--symbols', nargs='*', help='Symbols to include (space-separated)')
    p.add_argument('--config', default='config/rules.yaml', help='Config file to read symbols from')
    p.add_argument('--title', default='回測成果簡報')
    p.add_argument('--out', default=os.path.join('slides', 'final_20pages.pptx'))
    return p.parse_args()


if __name__ == '__main__':
    raise SystemExit(build(parse_args()))
